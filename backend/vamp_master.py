#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Offline VAMP core — ingestion + deep text extraction only.

This variant is intentionally stripped of NWUScorer / agent_app / Playwright
agent dependencies so that it can run as a standalone backend for the
Tkinter offline GUI:

    frontend/offline_app/offline_app_gui_llm_csv.py

It exposes two main functions:

    - extract_text_for(path: Path, size_limit: int = 200_000) -> str
    - ingest_paths(evidence_root: Path,
                   start: Optional[datetime.datetime],
                   end: Optional[datetime.datetime]) -> List[Artefact]

No NWU Brain deterministic scoring is performed here — that is handled
separately by the LLM layer (contextual_scorer.py) or by a full online
VAMP stack.
"""

from __future__ import annotations

import datetime
import os
import uuid
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional

# -------------------------
# Optional OCR support
# -------------------------

OCR_AVAILABLE = False
_OCR_ERROR = None
try:
    import pytesseract  # type: ignore
    from PIL import Image  # type: ignore
    from pdf2image import convert_from_path  # type: ignore

    OCR_AVAILABLE = True
except ImportError as e:  # pragma: no cover - optional deps
    _OCR_ERROR = str(e)

# -------------------------
# Utilities
# -------------------------


def sha1_file(path: Path) -> str:
    import hashlib

    h = hashlib.sha1()
    with path.open("rb") as fp:
        for chunk in iter(lambda: fp.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def guess_relpath(root: Path, p: Path) -> str:
    try:
        return str(p.relative_to(root))
    except Exception:
        return p.name


def generate_run_id() -> str:
    """Create a stable-ish identifier for a single scan run."""

    stamp = datetime.datetime.now().strftime("%Y%m%d%H%M%S")
    return f"run-{stamp}-{uuid.uuid4().hex[:8]}"


# -------------------------
# Text extraction (Deep Read)
# -------------------------


@dataclass
class ExtractionResult:
    extracted_text: str
    extract_status: str
    extract_error: Optional[str] = None


def _bytes_decode_guess(raw: bytes) -> str:
    """Best-effort bytes→str with chardet fallback."""
    try:
        import chardet  # type: ignore

        enc = chardet.detect(raw).get("encoding") or "utf-8"
        return raw.decode(enc, errors="ignore")
    except Exception:
        return raw.decode("utf-8", errors="ignore")


def txt_from_pdf(path: Path) -> str:
    """Extract text & table-ish content; robust to mildly corrupted PDFs."""
    text = ""
    # Attempt to validate/repair
    try:
        import pikepdf  # type: ignore

        with pikepdf.open(str(path)):
            pass
    except Exception:
        pass
    # pdfminer text layer
    try:
        from pdfminer.high_level import extract_text  # type: ignore

        text = extract_text(str(path)) or ""
    except Exception:
        text = ""
    # pdfplumber tables
    try:
        import pdfplumber  # type: ignore

        with pdfplumber.open(str(path)) as pdf:
            rows: List[str] = []
            for pg in pdf.pages:
                try:
                    tables = pg.extract_tables() or []
                    for tbl in tables:
                        rows.extend(
                            [" ".join([c or "" for c in r]) for r in tbl]
                        )
                except Exception:
                    pass
            if rows:
                text += "\n" + "\n".join(rows)
    except Exception:
        pass

    # OCR fallback for scanned PDFs (if text extraction failed)
    if (not text or len(text.strip()) < 50) and OCR_AVAILABLE:
        try:
            print(f"[OCR] Extracting text from scanned PDF: {path.name}")
            images = convert_from_path(str(path), dpi=300)
            ocr_text: List[str] = []
            for i, img in enumerate(images):
                page_text = pytesseract.image_to_string(
                    img, config="--psm 6"
                )
                if page_text.strip():
                    ocr_text.append(f"[Page {i+1}]\n{page_text}")
            if ocr_text:
                text = "\n\n".join(ocr_text)
                print(
                    f"[OCR] Extracted {len(text)} characters from "
                    f"{len(images)} pages"
                )
            else:
                print(f"[OCR] No text found in {path.name}")
        except Exception as e:  # pragma: no cover - best-effort logging
            print(f"[OCR] Failed for {path.name}: {e}")
    elif not text or len(text.strip()) < 50:
        # No OCR available and no text extracted
        print(
            f"[WARNING] No text extracted from {path.name} "
            f"(OCR not available)"
        )

    return text


def txt_from_docx(path: Path) -> str:
    try:
        import docx  # type: ignore

        d = docx.Document(str(path))
        return "\n".join(p.text for p in d.paragraphs)
    except Exception:
        return ""


def txt_from_xlsx(path: Path) -> str:
    try:
        import openpyxl  # type: ignore

        wb = openpyxl.load_workbook(
            str(path), data_only=True, read_only=True
        )
        out: List[str] = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                out.append(
                    " ".join("" if c is None else str(c) for c in row)
                )
        return "\n".join(out)
    except Exception:
        return ""


def txt_from_pptx(path: Path) -> str:
    try:
        from pptx import Presentation  # type: ignore

        prs = Presentation(str(path))
        out: List[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    out.append(shape.text or "")
        return "\n".join(out)
    except Exception:
        return ""


def _is_image(ext: str) -> bool:
    return ext in {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp"}


def extract_text_for(path: Path, size_limit: int = 200_000) -> ExtractionResult:
    """Deep Read extraction that always reports a status and error context."""

    name = path.name
    ext = path.suffix.lower()
    text = ""
    status = "ok"
    error: Optional[str] = None

    try:
        if ext == ".pdf":
            text = txt_from_pdf(path)
        elif ext == ".docx":
            text = txt_from_docx(path)
        elif ext in {".xlsx", ".xls", ".xlsm"}:
            text = txt_from_xlsx(path)
            if not text.strip():
                status = "empty_sheet"
        elif ext == ".pptx":
            text = txt_from_pptx(path)
        elif ext in {".txt", ".md", ".csv", ".log"}:
            text = _bytes_decode_guess(path.read_bytes())
        elif _is_image(ext):
            if not OCR_AVAILABLE:
                status = "image_no_ocr"
                error = _OCR_ERROR or "OCR dependencies not available"
                text = "IMAGE_EVIDENCE_REQUIRES_MANUAL_REVIEW"
            else:
                try:
                    img = Image.open(path)  # type: ignore[arg-type]
                    text = pytesseract.image_to_string(img, config="--psm 6")  # type: ignore[name-defined]
                    if not text.strip():
                        status = "failed"
                        error = "OCR returned no text"
                except Exception as ocr_exc:  # pragma: no cover - best-effort logging
                    status = "failed"
                    error = str(ocr_exc)
        elif ext == ".zip":
            try:
                with zipfile.ZipFile(str(path), "r") as zf:
                    inner = zf.namelist()[:80]
                    text = "ZIP " + " | ".join(inner)
                    status = "unsupported"
            except Exception as zip_exc:
                status = "failed"
                error = str(zip_exc)
        else:
            status = "unsupported"
            text = name
    except Exception as exc:  # pragma: no cover - defensive
        text = ""
        status = "failed"
        error = str(exc)

    text = text or ""
    if len(text) > size_limit:
        text = text[:size_limit]

    if status == "ok" and not text.strip():
        status = "failed"
        error = error or "no text extracted"

    return ExtractionResult(extracted_text=text, extract_status=status, extract_error=error)


# -------------------------
# Ingestion
# -------------------------

SKIP_DIRS = {"_out", "_final", "_logs", ".git", "__pycache__"}


@dataclass
class Artefact:
    path: Path
    relpath: str
    size: int
    mtime: float
    sha1: str


def in_window(
    ts: float,
    start: Optional[datetime.datetime],
    end: Optional[datetime.datetime],
) -> bool:
    """Check file mtime (POSIX seconds) falls inside window; permissive if not configured."""
    if not start or not end:
        return True
    mt = datetime.datetime.fromtimestamp(ts)
    return start <= mt < end


def ingest_paths(
    evidence_root: Path,
    start: Optional[datetime.datetime],
    end: Optional[datetime.datetime],
) -> List[Artefact]:
    """Walk the evidence_root and return unique Artefact entries within a time window.

    - Duplicates are removed based on SHA1 hash.
    - Common technical/auxiliary directories are skipped.
    """
    items: List[Artefact] = []
    seen: set[str] = set()
    for root, dirs, files in os.walk(evidence_root):
        # prune skip dirs in-place for performance
        dirs[:] = [d for d in dirs if d not in SKIP_DIRS]
        for fn in files:
            p = Path(root) / fn
            try:
                st = p.stat()
            except Exception:
                continue
            if not in_window(st.st_mtime, start, end):
                continue
            try:
                h = sha1_file(p)
            except Exception:
                h = f"ERR::{p.as_posix()}::{st.st_size}"
            if h in seen:
                continue
            seen.add(h)
            items.append(
                Artefact(
                    path=p,
                    relpath=guess_relpath(evidence_root, p),
                    size=st.st_size,
                    mtime=st.st_mtime,
                    sha1=h,
                )
            )
    return items
